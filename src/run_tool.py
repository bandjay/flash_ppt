# -*- coding: utf-8 -*-
"""
Created on Fri May 26 12:52:54 2017
@author: Jay Bandlamudi
"""
import os
import time
import argparse
from pptx import Presentation
from multiprocessing.dummy import Pool as ThreadPool
#from process_annotated_vcf import process_annotated_vcf             #### importing process_vcf.py to processs vcf file and to obtain a pandas dataframe 
#from generate_ppt_content import generate_ppt_content               #### importing ppt_content.py to generate report content from annotated vcf and web
#from fill_ppt_with_content import fill_ppt_with_content             #### importing ppt_fill.py to fille the empty presentation with content
#sys.dont_write_bytecode=True                                        #### to avoid creating .pyc file

def main(input_vcf_path,report_name,output_report_path,database_path): 
    input_vcf_path,report_name,output_report_path,database_path=input_vcf_path,report_name,output_report_path,database_path
    start = time.time()    
    src_path=os.getcwd()
    #home="C:/flash_ppt"
    home=os.path.split(src_path)[0]
    ppt_template_path= os.path.join(home,"ppt_template") 
    #input_vcf_path="C:/flash_ppt/sample_vcfs/sample_vcf"
    #output_report_path="C:/flash_ppt"
    #database_path="C:/flash_ppt"    
    """ 
    calling process_vcf to process the VCF file and generate a pandas dataframe 
    """ 
    #src_path="C:/flash_ppt/src"
    os.chdir(src_path)
    from process_annotated_vcf import process_annotated_vcf             #### importing process_vcf.py to processs vcf file and to obtain a pandas dataframe 
    pr=process_annotated_vcf(input_vcf_path)
    variants_dataframe=pr.process() 
       
    """
    ########################################################
    ########### Creating pptx ##############################
    ########################################################
    """   
    print " ===>  presentation generation is  in progress \n"        
    """ 
    using mayo custom template : template can be changed in the ppt_template directory 
    Note : template needs to contain 62 slides in order to facilitate 30 variants, refer existing template
    """
    os.chdir(ppt_template_path)    
    prs = Presentation("template_mayo.pptx")      
    os.chdir(src_path)
    from generate_ppt_content import generate_ppt_content               #### importing ppt_content.py to generate report content from annotated vcf and web
    from fill_ppt_with_content import fill_ppt_with_content             #### importing ppt_fill.py to fille the empty presentation with content
    """
    filling subject info slide with variant table
    """
    content_obj=generate_ppt_content(None)
    subject_slide_content=content_obj.subject_slide_content(variants_dataframe)
    pptobj=fill_ppt_with_content(None)
    pptobj.subject_slide(prs.slides[1],subject_slide_content)
                 
    """
    Dropping duplicate variants by GENE_SYM
    At this step creating new column "GENE_SYM" in the variants dataframe by parsing the 'INFO' field
 
    Info_rec=[variants_dataframe["INFO"][i] for i in range(0,variants_dataframe.shape[0])]
    Gene_sym_arr=[]
    for i in Info_rec:
        gene_dict={}
        gene=str(i).split(';')        
        #url=[]
        #non_url=[]
        for g in gene:
            if g.__contains__("href") or g.__contains__("www"):
               # url.append(g)
               continue                
            else:
                if g.__contains__("="):
                    #non_url.append(g)
                    k,v=g.split("=")
                    gene_dict[k]=v 
        if 'gene37p13.gene' in gene_dict.keys():
            Gene_sym_arr.append(gene_dict['gene37p13.gene'])
        else:
            raise Exception("GENE SYMBOL NOT FOUND IN NCBI / PLEASE CHECK NCBI CATALOG")
    variants_dataframe['GENE_SYM']=Gene_sym_arr      
    
    variants_dataframe=variants_dataframe.drop_duplicates('GENE_SYM', keep="last").reset_index(drop=True)
    """
    """
    slide_arr contains slide index of ppt template , for example first slide is for mayo information second slide is for
    subject and all variants information.From third slide onwards conatins gene information with 2 slides for each gene
    """    
    slide_arr=[[i+2,i+3] for i in range(2*variants_dataframe.shape[0]) if i%2==0] # in order to have 2 slides for a single gene
    
    """ 
    ##########################################################
    #################  calling ppt_content.py ################
    ##########################################################
    """    
    
    all_content_dict={}
    def get_content(data_record):
        pc=generate_ppt_content(database_path)
        return pc.generate_content(data_record)
    
    params=variants_dataframe.values.tolist()
    ''' multi-threaded web search for genes '''
    pool = ThreadPool(8) 
    params=variants_dataframe.values.tolist()
    all_content_dict= pool.map(get_content, params) 
    
    print " ===>  all slides content generated \n"
    
    print " ===>  Filling deck with content  \n"  
    '''
    i=1
    for p in params:
        print i
        get_content(p)
        i=i+1
    '''    
    """ 
    ##########################################################
    #################  calling ppt_fill.py ###################
    ##########################################################
    """
    for r in range(variants_dataframe.shape[0]):
        """ create a blank slide """      
        slide_num=slide_arr[r]
        """ calling fucntion to fill slide with contents """
        slide1=prs.slides[slide_num[0]]
        slide2=prs.slides[slide_num[1]]
   
        ppf=fill_ppt_with_content(all_content_dict[r])
        ppf.slideone(slide1)        
        ppf.slidetwo(slide2)
     
    """
    Template has 60 empty slides for each chromosome removing empty slides based on df size
    """ 
    
    for i in range(61,2*variants_dataframe.shape[0]+1,-1) : 
       rId = prs.slides._sldIdLst[i].rId
       prs.part.drop_rel(rId)
       del prs.slides._sldIdLst[i]
       
    '''
    saving presentation to output directory
    ''' 
    os.chdir(output_report_path)
    prs.save(str("sample_report")+".pptx")
    
    print " ===>  Presentation is in the directory : ",output_report_path,"\n"
    end = time.time()
    print " ===>  Time elapsed : ",end-start, "seconds \n"
    print "********************************************************"
    print "******************       END       *********************"
    print "********************************************************"
        
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Tool to generate automated report for variants')
    parser.add_argument('-i','--input', help='input vcf file name with allowed formats .vcf or .vcf.gz', required=True)
    parser.add_argument('-o','--output', help='output path to save the report', required=True)
    parser.add_argument('-n','--report_name', help='file name you prefer for the report', required=True)
    parser.add_argument('-d','--database', help='path to database, which stores gene details', required=True)
    
    args = vars(parser.parse_args())
    input_vcf_path=args['input']
    report_path=args['report_name']
    output_report_path=args['output']
    database_path=args['database']
    
    print "********************************************************"
    print "******************   STARTING FLASH ********************"
    print "********************************************************"
    
    main(input_vcf_path,report_path,output_report_path,database_path)       




