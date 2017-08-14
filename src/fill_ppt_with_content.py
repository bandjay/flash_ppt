
"""
Created on Fri Jun 23 09:51:28 2017
@author: Jay
"""

""" class to fill ppt with contents """
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

class fill_ppt_with_content:
    def __init__(self,content_dict):
        self.content_dict=content_dict 

    
    """ function to fill the subject slide with variants table """    
    def subject_slide(self,subslide,content):
        
        """ variants table """
        rows = 1+len(content)
        cols = 3
        left = Inches(1)
        top = Inches(3)
        width = Inches(4.0)
        height = Inches(0.8)
        
        table = subslide.shapes.add_table(rows, cols, left, top, width, height).table
                                         
        ''' font,column width adjustment in case of many variants '''
        font_size=Pt(12)
        col_width=Inches(2.5)
        if (rows>5 and  rows<=10):
            font_size=Pt(10)
            col_width=Inches(2)
        elif (rows>10 and rows<=20):
            font_size=Pt(8)   
            col_width=Inches(1.5)
        elif (rows>20):
            font_size=Pt(5)
            col_width=Inches(1)
        
        # set column widths
        table.columns[0].width = table.columns[1].width = table.columns[2].width = col_width
        
        # write column headings
        table.cell(0, 0).text = 'Gene Symbol'
        table.cell(0, 1).text = 'Transcript'
        table.cell(0, 2).text = 'Protein'
        #table.cell(0, 3).text = 'In silico SIFT/Polyphen'
        
        
        ''' Loop to fill table contents '''      
        
        for r in range(len(content)):
            val_arr=content[r]            
            for c in range(cols):
                table.cell(r+1,c).text=val_arr[c]       
        
        # change text font in table
        def iter_cells(table):
            for row in table.rows:
                for cell in row.cells:
                    yield cell

        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = font_size

    
    """ function to fill first slide of a gene"""    
    def slideone(self,slide1):  
        
        ''' adding a text box for Inheritance'''
        left = Inches(8.5)
        top = Inches(0.01)
        width = Inches(2)
        height = Inches(2)
        txBox = slide1.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.auto_size=True
        tf.word_wrap=True
        
        ''' title '''
        p = tf.add_paragraph()
        p.font.bold = True
        p.font.size=Pt(24)
        p.text=self.content_dict['Inheritance']
        p.font.color.rgb = RGBColor(255, 0, 0)
        
        
        ''' adding a text box '''
        left = Inches(0.01)
        top = Inches(0.01)
        width = Inches(10)
        height = Inches(5)
        txBox = slide1.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.auto_size=True
        tf.word_wrap=True
        
        ''' title '''
        p = tf.add_paragraph()
        p.text = "Variants in Genes Related to Phenotype"
        p.font.bold = True
        p.font.size=Pt(28) 
        p.font.color.rgb = RGBColor(0, 0, 128)
        
        ''' line 1 :GENE deatils'''
        p = tf.add_paragraph()
        p.text = self.content_dict['Gene']
        p.font.italic = True
        p.font.bold = True
        p.font.size=Pt(18)
        p.font.color.rgb = RGBColor(128, 0, 0)
        
        ''' line 2 : Disease'''
        """    
        Disease="Alzheimer disease type 3, MIM: 607822 & Alzheimer\
        disease type 3 with spastic paraparesis and unusual plaques,\
        MIM: 607822 & Alzheimer disease type 3 with spastic paraparesis\
        and apraxia, MIM: 607822 & Dementia frontotemporal, MIM: 600274\
        & Pick disease, MIM: 172700 & Cardiomyopathy dilated 1U, MIM: 613\
        694 & Acne inversa familial 3, MIM: 613737"
        
        disease_processed,MIM_ref_ids=process_disease(Disease)
        """ 
        
        def process_disease(Disease):
            dis_ARR=Disease.split("&")
            MIM_dict={}
            for di in dis_ARR:
                MIM=di.split("MIM:")[1].strip()
                Dis=di.split(",")[0]
                if MIM in MIM_dict.keys():
                    prev=MIM_dict[MIM]
                    MIM_dict[MIM]=prev+" & "+Dis
                else:
                    MIM_dict[MIM]=Dis
            process_ARR=[]
            for k in MIM_dict.keys():
                process_ARR.append(MIM_dict[k]+ ", MIM : " +k)
            MIM_ids=MIM_dict.keys()
            return (process_ARR,MIM_ids)
        disease_processed,MIM_ref_ids=process_disease(self.content_dict['Disease'])        
        
        p = tf.add_paragraph()
        p.text = "Disease and MIM List => "
        p.font.italic = True
        p.font.size=Pt(15)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0,0, 128)
        
        for dis in  disease_processed:  
            p = tf.add_paragraph()
            p.text = dis.strip()
            p.font.italic = True
            p.font.bold = True
            p.font.size=Pt(13)
            p.font.color.rgb = RGBColor(0,0, 128)
        
        ''' line 3 : Exac& gnomAD'''
        p = tf.add_paragraph()
        p.text = self.content_dict['Exac']
        p.font.bold = True
        p.font.size=Pt(15)
        p.font.color.rgb = RGBColor(0, 0, 128)    
        
        
        ''' line 4 : Clincar & HGMD '''
        p = tf.add_paragraph()
        p.text = self.content_dict['Clinvar']
        p.font.bold = True
        p.font.size=Pt(15)
        p.font.color.rgb = RGBColor(0, 0, 128)
        
        ''' line 5 : In Silico '''
        p = tf.add_paragraph()
        p.text = self.content_dict['In_silico']
        p.font.bold = True
        p.font.size=Pt(15)
        p.font.color.rgb = RGBColor(0, 0, 128)
        
        ''' line 6 : Location '''
        p = tf.add_paragraph()
        p.text = self.content_dict['Loc']
        p.font.bold = True
        p.font.size=Pt(15)
        p.font.color.rgb = RGBColor(0, 0, 128)
        
        
        ''' line 7 : Entrez Summary '''               
        p = tf.add_paragraph()
        p.text = self.content_dict['entrez_summary']
        p.font.italic = True
        p.font.size=Pt(13)
        p.font.color.rgb = RGBColor(0, 0, 0)
              
        
        ''' line 8 : Uniprot Summary                
        p = tf.add_paragraph()
        p.text = self.content_dict['uniprot_summary']
        p.font.italic = True
        p.font.size=Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)
        #p.font.size=Pt(12)
        '''
        
        if self.content_dict['exac_table']!="NA":
            ''' exac table on the right bottom corner'''        
            rows = 5
            cols = 4
            left = Inches(6)
            top = Inches(6)
            width = Inches(4.0)
            height = Inches(0.8)
            
            table = slide1.shapes.add_table(rows, cols, left, top, width, height).table
            
            # set column widths
            table.columns[0].width = Inches(1)
            table.columns[1].width = Inches(1)
            table.columns[2].width = Inches(1)
            table.columns[3].width = Inches(1)
            
            # write column headings
            table.cell(0, 0).text = 'Constraint from Exac'
            table.cell(0, 1).text = 'Expected no. variants'
            table.cell(0, 2).text = 'Observed no. variants'
            table.cell(0, 3).text = 'Constraint Metric'
            # write row headings
            table.cell(1, 0).text = 'Synonymous'
            table.cell(2, 0).text = 'Missense'
            table.cell(3, 0).text = 'LoF'
            table.cell(4, 0).text = 'CNV'
            
            # assigning table cellls with values
            #try:
            record=self.content_dict['exac_table'].split("\t")
            #record=exac_table.split("\t")
            table.cell(1, 1).text=str(record[0])
            table.cell(1, 2).text=str(record[1])
            table.cell(1, 3).text=str(record[2])
            table.cell(2, 1).text=str(record[3])
            table.cell(2, 2).text=str(record[4])
            table.cell(2, 3).text=str(record[5])
            table.cell(3, 1).text=str(record[6])
            table.cell(3, 2).text=str(record[7])
            table.cell(3, 3).text=str(record[8])
            table.cell(4, 1).text=str(record[9])
            table.cell(4, 2).text=str(record[10])
            table.cell(4, 3).text=str(record[11])
            #except:
            #    pass
                # change text font in table
            def iter_cells(table):
                    for row in table.rows:
                        for cell in row.cells:
                            yield cell
    
            for cell in iter_cells(table):
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)

        ''' OMIM URL reference for diseases '''
        left = Inches(3)
        top = Inches(5.75)
        width = height = Inches(2.5)
        txBox1 = slide1.shapes.add_textbox(left, top, width, height)
        tf1 = txBox1.text_frame
        p = tf1.add_paragraph()
        p.text="OMIM LINKS"
        p.font.bold = True
        p.font.size=Pt(11)
        p.font.color.rgb = RGBColor(225,0,0)
        for MI in MIM_ref_ids:
            p = tf1.add_paragraph()
            p.font.size=Pt(10)       
            run = p.add_run()
            run.text = str ('MIM :' + MI)
            hlink = run.hyperlink
            hlink.address = str("https://www.omim.org/entry/"+MI)
        
        ''' URL References on left bottom corner'''
        left = Inches(0.01)
        top = Inches(5.75)
        width = height = Inches(2.5)
        txBox1 = slide1.shapes.add_textbox(left, top, width, height)
        tf1 = txBox1.text_frame
        p = tf1.add_paragraph()
        p.text="WEB LINKS [right click]"
        p.font.bold = True
        p.font.size=Pt(11)
        p.font.color.rgb = RGBColor(225,0,0)
        
        '''
        p = tf1.add_paragraph()
        p.font.size=Pt(10)
        run = p.add_run()
        run.text = 'ExAC'
        hlink = run.hyperlink
        hlink.address = self.content_dict['exac_url']
        '''
        if self.content_dict['omim_url']!="NA":
            p = tf1.add_paragraph()
            p.font.size=Pt(10)       
            run = p.add_run()
            run.text = 'OMIM'
            hlink = run.hyperlink
            hlink.address = self.content_dict['omim_url']        
        if self.content_dict['ncbi_url']!="NA":
            p = tf1.add_paragraph()
            p.font.size=Pt(10)       
            run = p.add_run()
            run.text = 'NCBI'
            hlink = run.hyperlink
            hlink.address = self.content_dict['ncbi_url']
            #pptx.action.Hyperlink.address=hlink.address
        if self.content_dict['genecards_url']!="NA":
            p = tf1.add_paragraph()
            p.font.size=Pt(10)
            run = p.add_run()
            run.text = 'GENECARDS'
            hlink = run.hyperlink
            hlink.address = self.content_dict['genecards_url']
        if self.content_dict['clinvar_url']!="NA":
            p = tf1.add_paragraph()
            p.font.size=Pt(10)       
            run = p.add_run()
            run.text = 'CLINVAR_PUBMED'
            hlink = run.hyperlink
            hlink.address = self.content_dict['clinvar_url']
        
        if self.content_dict['hgmd_pub_url']!="NA":
            p = tf1.add_paragraph()
            p.font.size=Pt(10)       
            run = p.add_run()
            run.text = 'HGMD_PUBMED'
            hlink = run.hyperlink
            hlink.address = self.content_dict['hgmd_pub_url']

                
        
        """
        if self.content_dict['uniprot_url']!="NA":
            p = tf1.add_paragraph()
            p.font.size=Pt(10)       
            run = p.add_run()
            run.text = 'UNIPROT'
            hlink = run.hyperlink
            hlink.address = self.content_dict['uniprot_url']
            
        
        p = tf1.add_paragraph()
        p.font.size=Pt(10)
        run = p.add_run()
        run.text = 'marrvel'
        hlink = run.hyperlink
        hlink.address = self.content_dict['marrvel_url']
        
        
        
        """

    
    def slidetwo(self,slide2):
        
        ''' adding a text box '''
        left = Inches(0.25)
        top = Inches(0.01)
        width = Inches(10)
        height = Inches(5)
        txBox = slide2.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.auto_size=True
        tf.word_wrap=True
        
        p = tf.add_paragraph()
        p.text = "Phenotype List"
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 128)
        p.font.size=Pt(26)
        
        
        def pheno_print(pheno):
            p = tf.add_paragraph()
            p.text = pheno
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.font.size=Pt(16)
        
        pheno_type_list=self.content_dict['phenotype_list']
        if len(pheno_type_list)>20 :
            single_box_length=20
        else:
            single_box_length=len(pheno_type_list)
        for p in range(0,single_box_length):
            pheno_print(pheno_type_list[p])
        
        if len(pheno_type_list)>40 :
            ''' adding a text box '''
            left = Inches(5)
            top = Inches(0.8)
            width = Inches(6)
            height = Inches(5)
            txBox = slide2.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.auto_size=True  
            tf.word_wrap=True
            for p in range(20,40):
                pheno_print(pheno_type_list[p])
            p = tf.add_paragraph()
            p.text = "There are more ... , refer NOTES below"
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 128 )
            p.font.size=Pt(16)
            
            notes_slide = slide2.notes_slide
            text_frame = notes_slide.notes_text_frame
            extra_pheno_list=str(" , ".join(pheno_type_list[40:]))
            text_frame.text = extra_pheno_list

            
        
            
        
        